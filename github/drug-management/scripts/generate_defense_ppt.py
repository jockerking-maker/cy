# -*- coding: utf-8 -*-
"""
Generate graduation defense PPT — 简约蓝色学术风
Reference: 51PPT / 计算机毕设答辩常见模板（蓝白渐变 + 几何装饰 + 学术排版）
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT = os.path.join(os.path.dirname(os.path.dirname(__file__)), "答辩PPT_医院药品管理系统.pptx")

# ── 简约蓝色学术风配色（参考主流毕设模板：主色≤3种） ─────────────
PRIMARY    = RGBColor(0x1F, 0x4E, 0x79)   # 深藏青 — 标题栏/表头
ACCENT     = RGBColor(0x2E, 0x75, 0xB6)   # 标准蓝 — 强调/边框
LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)   # 浅蓝 — 卡片背景
SKY        = RGBColor(0xE8, 0xF4, 0xFC)   # 天蓝 — 封面/装饰
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE  = RGBColor(0xF8, 0xFA, 0xFD)
DARK_TEXT  = RGBColor(0x2C, 0x3E, 0x50)
MID_TEXT   = RGBColor(0x5A, 0x6A, 0x7A)
GRAY       = RGBColor(0x95, 0xA5, 0xA6)
BORDER     = RGBColor(0xBD, 0xD7, 0xEE)
ROW_ALT    = RGBColor(0xED, 0xF4, 0xFA)
SUCCESS    = RGBColor(0x27, 0xAE, 0x60)
GOLD       = RGBColor(0xF3, 0x9C, 0x12)   # 步骤序号点缀

THESIS_TITLE = "基于 Spring Boot 的医院药品管理系统的设计与实现"
STUDENT = "陈煜"
STUDENT_ID = "2022061008000403"
TOTAL = 17

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def solid(shape, color, transparency=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency is not None:
        shape.fill.transparency = transparency
    shape.line.fill.background()


def set_bg(slide, color=OFF_WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color, line=None, lw=1):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    return sh


def rrect(slide, l, t, w, h, fill=WHITE, line=BORDER, lw=1.2):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    return sh


def tb(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))


def p1(tf, text, size=18, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT, sa=8, sb=0):
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Microsoft YaHei"
    p.alignment = align
    p.space_after = Pt(sa)
    if sb:
        p.space_before = Pt(sb)
    return p


def pn(tf, text, **kw):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(kw.get("size", 18))
    p.font.bold = kw.get("bold", False)
    p.font.color.rgb = kw.get("color", DARK_TEXT)
    p.font.name = "Microsoft YaHei"
    p.alignment = kw.get("align", PP_ALIGN.LEFT)
    p.space_after = Pt(kw.get("sa", 8))
    if kw.get("sb"):
        p.space_before = Pt(kw["sb"])
    return p


def cover_decor(slide):
    """封面几何装饰 — 仿主流毕设模板"""
    set_bg(slide, WHITE)
    # 顶部浅蓝渐变块
    rect(slide, 0, 0, 13.333, 4.2, SKY)
    rect(slide, 0, 3.8, 13.333, 0.5, LIGHT_BLUE)
    # 右侧大圆
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-0.8), Inches(5.5), Inches(5.5))
    solid(c, LIGHT_BLUE, 0.35)
    # 左侧小圆
    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.2), Inches(5.5), Inches(3.5), Inches(3.5))
    solid(c2, SKY, 0.5)
    # 几何线条
    rect(slide, 0.5, 0.4, 2.5, 0.06, ACCENT)
    rect(slide, 0.5, 0.55, 0.06, 1.2, ACCENT)
    # 底部装饰点阵
    for i in range(8):
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8 + i * 0.35), Inches(6.85), Inches(0.08), Inches(0.08))
        solid(d, BORDER if i % 2 else ACCENT)


def slide_footer(slide, num):
    rect(slide, 0, 7.02, 13.333, 0.04, BORDER)
    box = rrect(slide, 11.6, 7.08, 1.35, 0.32, fill=SKY, line=BORDER, lw=0.8)
    t = tb(slide, 11.6, 7.1, 1.35, 0.3)
    p1(t.text_frame, f"{num:02d} / {TOTAL:02d}", size=11, color=ACCENT, align=PP_ALIGN.CENTER, sa=0)
    t2 = tb(slide, 0.5, 7.1, 10, 0.3)
    p1(t2.text_frame, "毕业设计答辩  ·  " + THESIS_TITLE[:26] + "…", size=9, color=GRAY, sa=0)


def content_slide(title, sec=None):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, OFF_WHITE)
    # 左侧蓝色竖条（模板标志元素）
    rect(slide, 0, 0, 0.18, 7.5, ACCENT)
    # 顶部标题栏
    rect(slide, 0.18, 0, 13.15, 1.12, PRIMARY)
    rect(slide, 0.18, 1.12, 13.15, 0.05, ACCENT)
    # 章节编号圆标
    if sec:
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.45), Inches(0.28), Inches(0.58), Inches(0.58))
        solid(circ, ACCENT)
        t = tb(slide, 0.45, 0.32, 0.58, 0.52)
        p1(t.text_frame, sec, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, sa=0)
        tx = 1.15
    else:
        tx = 0.55
    t = tb(slide, tx, 0.25, 11.5, 0.65)
    p1(t.text_frame, title, size=26, bold=True, color=WHITE, sa=0)
    # 内容白卡片
    rrect(slide, 0.45, 1.35, 12.45, 5.55, fill=WHITE, line=BORDER)
    return slide


def card_panel(slide, l, top, w, h, title, items, hc=PRIMARY):
    rrect(slide, l, top, w, h, fill=LIGHT_BLUE, line=BORDER)
    rect(slide, l, top, w, 0.07, ACCENT)
    title_box = tb(slide, l + 0.2, top + 0.15, w - 0.4, 0.42)
    p1(title_box.text_frame, title, size=19, bold=True, color=hc, sa=0)
    b = tb(slide, l + 0.25, top + 0.62, w - 0.45, h - 0.8)
    tf = b.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        txt = "●  " + item
        if i == 0:
            p1(tf, txt, size=15, color=DARK_TEXT, sa=10)
        else:
            pn(tf, txt, size=15, sa=10)


def styled_table(slide, headers, rows, l=0.65, top=1.55, w=12.05, col_widths=None):
    nrows, ncols = len(rows) + 1, len(headers)
    rh = min(0.5, 4.6 / nrows)
    ts = slide.shapes.add_table(nrows, ncols, Inches(l), Inches(top), Inches(w), Inches(rh * nrows))
    table = ts.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw)
    for j, h in enumerate(headers):
        c = table.cell(0, j)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = PRIMARY
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in c.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
            p.font.name = "Microsoft YaHei"
            p.alignment = PP_ALIGN.CENTER
    for i, row in enumerate(rows, 1):
        bg = ROW_ALT if i % 2 == 0 else WHITE
        for j, val in enumerate(row):
            c = table.cell(i, j)
            c.text = val
            c.fill.solid()
            c.fill.fore_color.rgb = bg
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.name = "Microsoft YaHei"
                p.font.color.rgb = DARK_TEXT
                p.alignment = PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT
                if "通过" in val:
                    p.font.color.rgb = SUCCESS
                    p.font.bold = True


def bullets(slide, items, l=0.75, top=1.55, w=11.8, size=16):
    b = tb(slide, l, top, w, 5.0)
    tf = b.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p1(tf, item, size=size, sa=11)
        else:
            pn(tf, item, size=size, sa=11)


def flow_steps(slide, steps, y=2.0):
    n = len(steps)
    bw = min(1.42, 11.2 / n - 0.15)
    gap = (11.2 - bw * n) / max(n - 1, 1)
    x = 1.0
    cols = [PRIMARY, ACCENT, PRIMARY, ACCENT, PRIMARY, ACCENT, ACCENT]
    for i, step in enumerate(steps):
        c = cols[i % len(cols)]
        rrect(slide, x, y, bw, 0.78, fill=c, line=None)
        t = tb(slide, x + 0.04, y + 0.16, bw - 0.08, 0.52)
        p1(t.text_frame, step, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, sa=0)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + bw / 2 - 0.14), Inches(y - 0.2), Inches(0.28), Inches(0.28))
        solid(circ, GOLD)
        nt = tb(slide, x + bw / 2 - 0.14, y - 0.17, 0.28, 0.25)
        p1(nt.text_frame, str(i + 1), size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER, sa=0)
        if i < n - 1:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + bw + 0.02), Inches(y + 0.3), Inches(gap - 0.04), Inches(0.16))
            solid(arr, BORDER)
        x += bw + gap


page = 1

# ── 1 封面 ────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
cover_decor(slide)
t = tb(slide, 0.8, 0.55, 11.7, 0.45)
p1(t.text_frame, "本科毕业设计答辩", size=15, color=ACCENT, align=PP_ALIGN.CENTER, sa=0)
t = tb(slide, 0.6, 1.1, 12.1, 2.2)
p1(t.text_frame, "基于 Spring Boot 的", size=34, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER, sa=6)
pn(t.text_frame, "医院药品管理系统", size=38, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER, sa=6)
pn(t.text_frame, "的设计与实现", size=26, color=ACCENT, align=PP_ALIGN.CENTER, sa=0)
rrect(slide, 3.0, 4.35, 7.3, 2.55, fill=WHITE, line=ACCENT, lw=1.5)
for i, (lb, val) in enumerate([
    ("答辩人", STUDENT), ("学  号", STUDENT_ID),
    ("专  业", "计算机科学与技术"), ("指导教师", "__________"), ("答辩日期", "2026 年 5 月"),
]):
    t = tb(slide, 3.3, 4.55 + i * 0.42, 6.7, 0.38)
    p1(t.text_frame, f"{lb}    {val}", size=17, color=DARK_TEXT if i == 0 else MID_TEXT, align=PP_ALIGN.CENTER, sa=0)
page += 1

# ── 2 目录 ────────────────────────────────────────────────────
slide = content_slide("汇报提纲", "目")
slide_footer(slide, page)
secs = [("01", "研究背景与意义"), ("02", "国内外研究现状"), ("03", "系统需求分析"),
        ("04", "系统总体设计"), ("05", "详细设计与实现"), ("06", "系统测试验证"), ("07", "总结与展望")]
pos = [(0.65, 1.55), (4.55, 1.55), (8.45, 1.55), (0.65, 3.55), (4.55, 3.55), (8.45, 3.55), (4.55, 5.55)]
for (num, title), (x, y) in zip(secs, pos):
    fill = LIGHT_BLUE if int(num) % 2 else WHITE
    rrect(slide, x, y, 3.6, 1.7, fill=fill, line=BORDER)
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.22), Inches(y + 0.52), Inches(0.62), Inches(0.62))
    solid(circ, PRIMARY if int(num) % 2 else ACCENT)
    nt = tb(slide, x + 0.22, y + 0.57, 0.62, 0.55)
    p1(nt.text_frame, num, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, sa=0)
    tt = tb(slide, x + 1.0, y + 0.62, 2.4, 0.6)
    p1(tt.text_frame, title, size=16, bold=True, color=DARK_TEXT, sa=0)
page += 1

# ── 3 背景 ────────────────────────────────────────────────────
slide = content_slide("研究背景与意义", "01")
slide_footer(slide, page)
card_panel(slide, 0.65, 1.55, 5.85, 5.1, "研究背景", [
    "医疗信息化进程加快，传统人工+纸质管理效率低",
    "药品管理涉及采购、入库、存储、出库、使用多环节",
    "业务流程复杂、数据量大、安全性要求高",
    "信息共享困难，易出现数据不一致和操作差错",
])
card_panel(slide, 6.75, 1.55, 5.85, 5.1, "研究意义", [
    "提高医院药品管理效率，降低管理成本",
    "实现药品全生命周期数字化管理",
    "保障临床用药安全，减少人为失误",
    "为医院信息化建设提供可复用的技术方案",
], hc=ACCENT)
page += 1

# ── 4 现状 ────────────────────────────────────────────────────
slide = content_slide("国内外研究现状", "02")
slide_footer(slide, page)
card_panel(slide, 0.65, 1.55, 5.85, 5.1, "国外研究现状", [
    "欧美医院普遍采用 HIS/ERP 集成化药品管理系统",
    "强调批次追溯、条码管理和自动化仓储",
    "注重与临床信息系统（CIS）的数据互通",
])
card_panel(slide, 6.75, 1.55, 5.85, 5.1, "国内研究现状", [
    "国内医院信息化起步较晚，但发展速度快",
    "Spring Boot + Vue 前后端分离架构应用广泛",
    "现有系统多侧重单一环节，缺乏全流程整合",
    "本系统面向中小型医院，覆盖采购到库存全链路",
], hc=ACCENT)
page += 1

# ── 5 需求 ────────────────────────────────────────────────────
slide = content_slide("系统需求分析", "03")
slide_footer(slide, page)
card_panel(slide, 0.65, 1.55, 5.85, 5.1, "功能需求 · 四大角色", [
    "管理员 — 用户管理、角色权限、系统公告",
    "采购员 — 采购订单管理、供应商管理",
    "采购审核员 — 采购单多级审核",
    "库管员 — 药品入库、出库、库存预警",
])
card_panel(slide, 6.75, 1.55, 5.85, 5.1, "非功能需求", [
    "易用性 — 界面简洁，操作流程清晰",
    "可靠性 — 事务机制保证数据一致性",
    "安全性 — JWT 认证 + RBAC 权限控制",
    "性能 — 页面响应 ≤ 3s，并发 ≥ 100 QPS",
    "可维护性 — 分层架构，模块化设计",
], hc=ACCENT)
page += 1

# ── 6 用例 ────────────────────────────────────────────────────
slide = content_slide("系统用例与业务流程", "03")
slide_footer(slide, page)
flows = [
    ("采购流程", "创建采购单 → 多级审核 → 采购入库 → 更新库存"),
    ("出库流程", "选择药品/批次 → 校验库存 → 扣减库存 → 生成记录"),
    ("预警流程", "检测库存/有效期 → 生成预警 → 指派处理 → 通知"),
    ("盘点流程", "创建盘点单 → 录入实盘 → 计算盈亏 → 调整库存"),
]
y = 1.6
for i, (title, desc) in enumerate(flows):
    rrect(slide, 0.65, y, 12.0, 1.0, fill=WHITE, line=ACCENT if i % 2 == 0 else BORDER)
    rect(slide, 0.65, y, 0.1, 1.0, PRIMARY if i % 2 == 0 else ACCENT)
    t = tb(slide, 0.9, y + 0.12, 2.2, 0.4)
    p1(t.text_frame, title, size=16, bold=True, color=PRIMARY, sa=0)
    t2 = tb(slide, 3.2, y + 0.22, 9.2, 0.5)
    p1(t2.text_frame, desc, size=15, color=MID_TEXT, sa=0)
    y += 1.15
t = tb(slide, 0.65, 6.3, 12, 0.4)
p1(t.text_frame, "用户登录后按 RBAC 模型分配菜单与操作权限（参见论文图 3-1）", size=13, color=GRAY, align=PP_ALIGN.CENTER, sa=0)
page += 1

# ── 7 技术选型 ────────────────────────────────────────────────
slide = content_slide("技术选型", "04")
slide_footer(slide, page)
styled_table(slide, ["层次", "技术 / 工具", "版本", "主要用途"], [
    ["后端框架", "Spring Boot", "3.2.5", "RESTful API · 依赖注入 · 事务管理"],
    ["ORM 框架", "MyBatis-Plus", "3.5.7", "CRUD 封装 · 条件构造器 · 分页"],
    ["数据库", "MySQL", "8.0", "关系型存储 · 索引优化"],
    ["安全认证", "JWT", "0.12.6", "无状态 Token 身份认证"],
    ["前端框架", "Vue 3", "3.4.0", "组件化开发 · 响应式 UI"],
    ["UI 组件", "Element Plus", "2.5.0", "表格 · 表单 · 对话框"],
    ["构建工具", "Vite", "—", "快速热更新开发"],
    ["数据可视化", "ECharts", "—", "采购 / 库存统计图表"],
], col_widths=[1.6, 2.2, 1.2, 7.05])
page += 1

# ── 8 架构 ────────────────────────────────────────────────────
slide = content_slide("系统架构设计", "04")
slide_footer(slide, page)
layers = [
    ("表现层", "Vue 3 · Element Plus · ECharts · Pinia · Axios"),
    ("接口层", "Spring MVC · RESTful API · JWT 认证 · 统一响应"),
    ("业务层", "Service · 事务管理 · 入库 / 出库 / 预警规则"),
    ("持久层", "MyBatis-Plus · Mapper · 参数化查询"),
    ("存储层", "MySQL 8.0 · 索引优化 · 事务管理"),
]
y = 1.55
cols = [PRIMARY, ACCENT, PRIMARY, ACCENT, PRIMARY]
for i, (name, desc) in enumerate(layers):
    c = cols[i]
    rrect(slide, 1.4, y, 10.5, 0.85, fill=WHITE, line=c, lw=2)
    rrect(slide, 1.55, y + 0.12, 1.3, 0.6, fill=c, line=None)
    t = tb(slide, 1.55, y + 0.18, 1.3, 0.5)
    p1(t.text_frame, name, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, sa=0)
    t2 = tb(slide, 3.05, y + 0.24, 8.6, 0.5)
    p1(t2.text_frame, desc, size=15, color=DARK_TEXT, sa=0)
    if i < 4:
        arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.3), Inches(y + 0.85), Inches(0.32), Inches(0.2))
        solid(arr, ACCENT)
    y += 1.02
t = tb(slide, 1.4, 6.65, 10.5, 0.35)
p1(t.text_frame, "B/S 架构 + 前后端分离  |  HTTP/REST  |  浏览器 + 应用服务器 + MySQL", size=13, color=GRAY, align=PP_ALIGN.CENTER, sa=0)
page += 1

# ── 9 模块 ────────────────────────────────────────────────────
slide = content_slide("功能模块设计", "04")
slide_footer(slide, page)
styled_table(slide, ["功能模块", "主要功能", "后端 Controller"], [
    ["系统管理", "用户管理 · 角色权限 · 系统公告", "SysUser / SysRole / SystemNotice"],
    ["基础数据", "药品信息 · 供应商 · 仓库管理", "DrugInfo / SupplierInfo / Warehouse"],
    ["采购管理", "采购单创建 · 多级审核 · 明细管理", "PurchaseOrderController"],
    ["库存管理", "库存查询 · 盘点 · 药品锁定", "Stock / StockCheck / DrugLock"],
    ["出入库", "采购入库 · 领用出库 · 批次管理", "DrugIn / DrugOutController"],
    ["预警管理", "低库存 · 高库存 · 临期预警", "StockWarningController"],
    ["报表统计", "采购 / 出入库 / 库存 / 财务分析", "DashboardController"],
], top=1.5, col_widths=[1.8, 5.5, 4.75])
page += 1

# ── 10 数据库 ─────────────────────────────────────────────────
slide = content_slide("数据库设计", "04")
slide_footer(slide, page)
for x, num, label in [(0.65, "19", "数据表"), (3.5, "4", "角色类型"), (6.35, "15+", "Controller"), (9.2, "8", "索引组")]:
    rrect(slide, x, 1.55, 2.5, 1.05, fill=LIGHT_BLUE, line=ACCENT)
    t = tb(slide, x, 1.62, 2.5, 0.5)
    p1(t.text_frame, num, size=30, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER, sa=0)
    t2 = tb(slide, x, 2.1, 2.5, 0.35)
    p1(t2.text_frame, label, size=13, color=MID_TEXT, align=PP_ALIGN.CENTER, sa=0)
bullets(slide, [
    "●  核心业务表：drug_info · drug_stock · purchase_order · drug_in · drug_out",
    "●  系统管理表：sys_user · sys_role · sys_menu · sys_role_menu",
    "●  辅助业务表：stock_warning · stock_check · audit_record · drug_lock",
    "●  设计原则：第三范式 · 逻辑删除 · 统一时间戳",
    "●  性能优化：drug_code · order_no 等关键字段建立唯一 / 联合索引",
], top=2.85, size=16)
page += 1

# ── 11 采购 ───────────────────────────────────────────────────
slide = content_slide("核心功能实现 — 采购管理", "05")
slide_footer(slide, page)
flow_steps(slide, ["创建采购单", "填写明细", "提交审核", "审核处理", "审核通过", "执行入库", "更新库存"])
rrect(slide, 0.65, 3.1, 12.0, 0.5, fill=SKY, line=None)
t = tb(slide, 0.8, 3.18, 11.7, 0.38)
p1(t.text_frame, "状态流转：待审核(0)  →  已审核(1)  →  已入库(2)  /  已取消(3)", size=15, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER, sa=0)
card_panel(slide, 0.65, 3.75, 5.85, 2.9, "关键技术", [
    "多级审核记录写入 audit_record 表",
    "审核通过后关联采购单执行入库",
    "@Transactional 保证入库与库存更新一致性",
])
card_panel(slide, 6.75, 3.75, 5.85, 2.9, "业务价值", [
    "规范采购审批流程，防止违规采购",
    "采购与库存数据自动联动",
    "完整操作日志，支持追溯审计",
], hc=ACCENT)
page += 1

# ── 12 库存预警 ───────────────────────────────────────────────
slide = content_slide("核心功能实现 — 库存与预警", "05")
slide_footer(slide, page)
card_panel(slide, 0.65, 1.55, 5.85, 5.1, "库存管理", [
    "药品 + 仓库 + 批次三维度精细化管理",
    "入库增加 quantity，记录生产日期 / 有效期",
    "出库校验可用库存，扣减 quantity",
    "lock_num 锁定机制防止超卖",
    "盘点对比系统数量与实盘，处理盈亏",
])
card_panel(slide, 6.75, 1.55, 5.85, 5.1, "库存预警", [
    "低库存：quantity < warning_num",
    "高库存：quantity > max_warning_num",
    "临期预警：距 expiry_date 不足阈值",
    "预警级别：一般 / 重要 / 紧急",
    "支持指派处理人及多渠道通知",
], hc=ACCENT)
page += 1

# ── 13 安全 ───────────────────────────────────────────────────
slide = content_slide("安全机制实现", "05")
slide_footer(slide, page)
sec = [("JWT 认证", "Token 机制\n未授权返回 401"), ("RBAC 权限", "角色-菜单\n三级授权"),
       ("密码加密", "MD5 加盐\n不存明文"), ("SQL 防护", "参数化查询\n防注入"),
       ("操作审计", "操作日志\n全程追溯"), ("跨域配置", "CorsConfig\n前后端分离")]
for idx, (title, desc) in enumerate(sec):
    col, row = idx % 3, idx // 3
    x, y = 0.65 + col * 4.1, 1.55 + row * 2.55
    rrect(slide, x, y, 3.75, 2.15, fill=WHITE, line=ACCENT if idx % 2 == 0 else BORDER)
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.22), Inches(y + 0.22), Inches(0.48), Inches(0.48))
    solid(circ, PRIMARY if idx % 2 == 0 else ACCENT)
    t = tb(slide, x + 0.85, y + 0.25, 2.7, 0.42)
    p1(t.text_frame, title, size=16, bold=True, color=PRIMARY, sa=0)
    t2 = tb(slide, x + 0.28, y + 0.88, 3.2, 1.0)
    for j, line in enumerate(desc.split("\n")):
        if j == 0:
            p1(t2.text_frame, line, size=14, color=MID_TEXT, sa=4)
        else:
            pn(t2.text_frame, line, size=14, color=MID_TEXT, sa=0)
page += 1

# ── 14 演示 ───────────────────────────────────────────────────
slide = content_slide("系统界面演示", "05")
slide_footer(slide, page)
styled_table(slide, ["步骤", "演示页面", "演示要点"], [
    ["①", "登录页", "JWT 登录认证 · 按角色跳转首页"],
    ["②", "系统仪表盘", "ECharts 图表 · 数据概览"],
    ["③", "药品管理", "CRUD · 上下架 · 条件查询"],
    ["④", "采购管理", "创建采购单 · 提交审核"],
    ["⑤", "采购审核", "审核通过 / 驳回 · 审核记录"],
    ["⑥", "药品入库", "关联采购单 · 批次入库 · 库存更新"],
    ["⑦", "库存预警", "预警列表 · 指派处理"],
    ["⑧", "报表统计", "采购 / 出入库趋势分析"],
], top=1.5, col_widths=[0.8, 2.5, 8.75])
t = tb(slide, 0.65, 6.5, 12, 0.35)
p1(t.text_frame, "建议现场演示：采购 → 审核 → 入库 完整业务链路", size=13, color=ACCENT, align=PP_ALIGN.CENTER, sa=0)
page += 1

# ── 15 测试 ───────────────────────────────────────────────────
slide = content_slide("系统测试与验证", "06")
slide_footer(slide, page)
styled_table(slide, ["测试类型", "测试内容", "测试结果"], [
    ["功能测试", "15 个核心用例（登录 / 采购 / 入库 / 出库 / 预警）", "全部通过 ✓"],
    ["性能测试", "10 并发用户 · 核心接口响应时间", "平均 < 200ms · QPS > 50"],
    ["安全测试", "认证 · 权限 · 密码 · SQL 注入防护", "全部通过 ✓"],
], top=1.55, col_widths=[1.8, 6.5, 3.75])
for x, title, val, sub in [
    (0.65, "测试环境", "Win10 + JDK17 + MySQL8.0", "Chrome 浏览器"),
    (4.55, "测试工具", "Postman + 浏览器", "接口 & 功能测试"),
    (8.45, "测试结论", "功能完整 · 运行稳定", "满足非功能需求"),
]:
    rrect(slide, x, 4.0, 3.6, 2.5, fill=LIGHT_BLUE, line=BORDER)
    t = tb(slide, x + 0.15, 4.12, 3.3, 0.38)
    p1(t.text_frame, title, size=15, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER, sa=0)
    t2 = tb(slide, x + 0.15, 4.65, 3.3, 0.48)
    p1(t2.text_frame, val, size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, sa=0)
    t3 = tb(slide, x + 0.15, 5.25, 3.3, 0.4)
    p1(t3.text_frame, sub, size=13, color=MID_TEXT, align=PP_ALIGN.CENTER, sa=0)
page += 1

# ── 16 总结 ───────────────────────────────────────────────────
slide = content_slide("总结与展望", "07")
slide_footer(slide, page)
card_panel(slide, 0.65, 1.55, 5.85, 5.1, "工作总结", [
    "完成需求分析、架构设计、数据库设计",
    "实现采购、入库、出库、库存、预警等核心功能",
    "采用 Spring Boot + Vue 3 前后端分离架构",
    "设计 19 张数据表，完善权限与安全机制",
    "通过功能、性能、安全三类测试验证",
])
card_panel(slide, 6.75, 1.55, 5.85, 5.1, "不足与展望", [
    "尚未对接 HIS 临床系统，数据互通待加强",
    "移动端适配不足，可开发小程序 / APP",
    "预警通知可接入短信 / 邮件实际通道",
    "可引入 Redis 缓存提升高并发性能",
    "可扩展条码 / RFID 自动化仓储管理",
], hc=ACCENT)
page += 1

# ── 17 致谢 ───────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
cover_decor(slide)
t = tb(slide, 0.8, 2.2, 11.7, 1.2)
p1(t.text_frame, "谢谢各位老师！", size=50, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER, sa=0)
pn(t.text_frame, "敬请批评指正", size=28, color=ACCENT, align=PP_ALIGN.CENTER, sb=14, sa=0)
rrect(slide, 5.15, 3.85, 3.1, 0.72, fill=PRIMARY, line=None)
t = tb(slide, 5.15, 3.93, 3.1, 0.6)
p1(t.text_frame, "Q  &  A", size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, sa=0)
t = tb(slide, 0.8, 5.2, 11.7, 0.45)
p1(t.text_frame, f"答辩人：{STUDENT}  |  学号：{STUDENT_ID}", size=16, color=MID_TEXT, align=PP_ALIGN.CENTER, sa=0)

prs.save(OUTPUT)
print(f"PPT generated: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
